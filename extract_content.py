"""Extract text and images from PPTX and DOCX files in Тема-1 folder."""
import os
import json
import base64
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from docx import Document
from PIL import Image
import io

BASE = Path(r"d:\presentation\Тема-1")
OUTPUT = Path(r"d:\presentation\extracted")
OUTPUT.mkdir(exist_ok=True)
IMAGES_DIR = OUTPUT / "images"
IMAGES_DIR.mkdir(exist_ok=True)

def extract_pptx(filepath, tag):
    prs = Presentation(filepath)
    slides_data = []
    img_counter = 0
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_info = {
            "slide": slide_num,
            "title": "",
            "content": [],
            "images": [],
            "notes": ""
        }
        # Get slide layout/title
        for shape in slide.shapes:
            if shape.has_text_frame:
                is_title = shape.name.lower().startswith("title")
                try:
                    ph = shape.placeholder_format
                    if ph and ph.idx == 0:
                        is_title = True
                except Exception:
                    pass
                if is_title:
                    slide_info["title"] = shape.text_frame.text.strip()
        # Get all text
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text and text != slide_info["title"]:
                        level = para.level
                        slide_info["content"].append({"level": level, "text": text})
            # Extract images
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                try:
                    image = shape.image
                    img_bytes = image.blob
                    ext = image.ext
                    img_name = f"{tag}_slide{slide_num}_img{img_counter}.{ext}"
                    img_path = IMAGES_DIR / img_name
                    with open(img_path, "wb") as f:
                        f.write(img_bytes)
                    # Also create base64
                    b64 = base64.b64encode(img_bytes).decode()
                    mime = f"image/{ext}" if ext != "jpg" else "image/jpeg"
                    slide_info["images"].append({
                        "file": img_name,
                        "base64": f"data:{mime};base64,{b64}",
                        "left": shape.left,
                        "top": shape.top,
                        "width": shape.width,
                        "height": shape.height
                    })
                    img_counter += 1
                except Exception as e:
                    print(f"  Image error slide {slide_num}: {e}")
        # Notes
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_info["notes"] = notes_text
        slides_data.append(slide_info)
    return slides_data

def extract_docx(filepath, tag):
    doc = Document(filepath)
    content = []
    img_counter = 0
    for elem in doc.element.body:
        tag_name = elem.tag.split('}')[-1] if '}' in elem.tag else elem.tag
        if tag_name == 'p':
            # Find paragraph in doc
            pass
    # Use paragraphs
    for para in doc.paragraphs:
        text = para.text.strip()
        style = para.style.name if para.style else ""
        if text:
            content.append({"style": style, "text": text})
    # Extract images
    images = []
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            try:
                img_part = rel.target_part
                img_bytes = img_part.blob
                content_type = img_part.content_type
                ext = content_type.split("/")[-1]
                if ext == "jpeg": ext = "jpg"
                img_name = f"{tag}_doc_img{img_counter}.{ext}"
                img_path = IMAGES_DIR / img_name
                with open(img_path, "wb") as f:
                    f.write(img_bytes)
                b64 = base64.b64encode(img_bytes).decode()
                images.append({
                    "file": img_name,
                    "base64": f"data:image/{ext};base64,{b64}"
                })
                img_counter += 1
            except Exception as e:
                print(f"  DOCX image error: {e}")
    return {"paragraphs": content, "images": images}

print("=== Extracting PPTX: Мултимедийна презентация ===")
pptx_main = extract_pptx(BASE / "Нива1-2_Тема1.1_Мултимедийна презентация.pptx", "main")
print(f"  {len(pptx_main)} slides")

print("=== Extracting PPTX: Интерактивна демонстрация ===")
pptx_demo = extract_pptx(BASE / "Нива1-2_Тема1.1_Интерактивна демонстрация.pptx", "demo")
print(f"  {len(pptx_demo)} slides")

print("=== Extracting DOCX: Лекционен материал ===")
docx_lecture = extract_docx(BASE / "Нива1-2_Тема1.1_Лекционен материал.docx", "lecture")
print(f"  {len(docx_lecture['paragraphs'])} paragraphs, {len(docx_lecture['images'])} images")

print("=== Extracting DOCX: Тестови въпроси ===")
docx_tests = extract_docx(BASE / "Нива1-2_Тема1.1_Тестови въпроси за самопроверка.docx", "tests")
print(f"  {len(docx_tests['paragraphs'])} paragraphs, {len(docx_tests['images'])} images")

# Save all to JSON (without base64 for readability)
def strip_b64(obj):
    if isinstance(obj, dict):
        return {k: ("[base64 omitted]" if k == "base64" else strip_b64(v)) for k, v in obj.items()}
    if isinstance(obj, list):
        return [strip_b64(i) for i in obj]
    return obj

with open(OUTPUT / "main_pptx.json", "w", encoding="utf-8") as f:
    json.dump(pptx_main, f, ensure_ascii=False, indent=2)
with open(OUTPUT / "demo_pptx.json", "w", encoding="utf-8") as f:
    json.dump(pptx_demo, f, ensure_ascii=False, indent=2)
with open(OUTPUT / "lecture_docx.json", "w", encoding="utf-8") as f:
    json.dump(docx_lecture["paragraphs"], f, ensure_ascii=False, indent=2)
with open(OUTPUT / "tests_docx.json", "w", encoding="utf-8") as f:
    json.dump(docx_tests, f, ensure_ascii=False, indent=2)

print("\n=== Summary (no base64) ===")
print(json.dumps(strip_b64(pptx_main), ensure_ascii=False, indent=2)[:8000])
